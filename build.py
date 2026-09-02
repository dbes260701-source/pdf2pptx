"""
build.py -- layout JSON (+ overrides) -> PPTX
usage: python build.py <workdir> <out.pptx> [--pages 1,2] [--font "맑은 고딕"]
Overrides: work/overrides/pN.json
  { "delete": ["l3","i2"],
    "set": {"t5": {"text": "corrected", "bold": true, "font_pt": 11}},
    "add": [ {type:"text",...}, {type:"line",...}, {type:"rect",...}, {type:"image", asset:"x.png", bbox:[..]} ] }
"""
import sys, os, json, glob
import numpy as np, cv2
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

def respath(*parts):
    base = getattr(sys, "_MEIPASS", os.path.dirname(os.path.abspath(__file__)))
    return os.path.join(base, *parts)

GLYPH_RATIO = 0.90   # korean glyph bbox height / font size (tuned by render compare)

WIN_FONTS = os.environ.get("WINDIR", r"C:\Windows") + r"\Fonts"

def _find_system_font(family, bold):
    """설치된 폰트 파일 경로 찾기 (배포 시 폰트 재배포를 피하기 위함)"""
    try:
        import glob as _g
        from PIL import ImageFont
        for p in _g.glob(os.path.join(WIN_FONTS, "*.ttf")) + _g.glob(os.path.join(WIN_FONTS, "*.ttc")):
            try:
                n = ImageFont.truetype(p, 20).getname()
            except Exception:
                continue
            if n[0].replace(" ", "").lower() == family.replace(" ", "").lower():
                if bold and n[1].lower() in ("bold",): return p
                if not bold and n[1].lower() in ("regular", "book"): return p
    except Exception:
        pass
    return None

_FONT_STATIC = {"맑은 고딕": (os.path.join(WIN_FONTS, "malgun.ttf"), os.path.join(WIN_FONTS, "malgunbd.ttf")),
                "Malgun Gothic": (os.path.join(WIN_FONTS, "malgun.ttf"), os.path.join(WIN_FONTS, "malgunbd.ttf")),
                "NanumGothic": (respath("fonts", "NanumGothic.ttf"), respath("fonts", "NanumGothicBold.ttf"))}
_font_path_cache = {}

def font_file(family, bold):
    key = (family, bold)
    if key in _font_path_cache: return _font_path_cache[key]
    cand = _FONT_STATIC.get(family, (None, None))[1 if bold else 0]
    if not (cand and os.path.exists(cand)):
        cand = _find_system_font(family, bold) or _find_system_font("Malgun Gothic", bold)                or os.path.join(WIN_FONTS, "malgunbd.ttf" if bold else "malgun.ttf")
    _font_path_cache[key] = cand
    return cand

_font_cache = {}
def measure_width(text, font, bold):
    """text width in px at 100px font size using PIL (for width-fitted sizing)"""
    from PIL import ImageFont
    path = font_file(font, bold)
    key = (path, 100)
    if key not in _font_cache: _font_cache[key] = ImageFont.truetype(path, 100)
    return max(1.0, _font_cache[key].getlength(text))

def rgb(hexs): hexs = hexs.lstrip("#"); return RGBColor(int(hexs[:2],16), int(hexs[2:4],16), int(hexs[4:],16))

def apply_overrides(layout, ov):
    els = {e["id"]: e for e in layout["elements"]}
    for d in ov.get("delete", []): els.pop(d, None)
    for k, v in ov.get("set", {}).items():
        if k in els: els[k].update(v)
    out = list(els.values())
    for k, a in enumerate(ov.get("add", [])):
        a.setdefault("id", f"a{k}"); out.append(a)
    layout["elements"] = out

def crop_image(work, layout, e, asset_dir):
    """crop from hi-res page; mask text inside crop with local bg; optional white->transparent"""
    i = layout["page"]; hi = cv2.imread(f"{work}/hires/p{i}.png")
    s = hi.shape[1] / layout["width"]
    x0,y0,x1,y1 = [int(round(v*s)) for v in e["bbox"]]
    crop = hi[y0:y1, x0:x1].copy()
    for t in layout["elements"]:
        if t["type"] != "text" or t.get("keep_in_image"): continue
        b = t["bbox"]
        ix = min(b[2], e["bbox"][2]) - max(b[0], e["bbox"][0]); iy = min(b[3], e["bbox"][3]) - max(b[1], e["bbox"][1])
        if ix > 0 and iy > 0:
            bx0,by0,bx1,by1 = [int(round(v*s)) for v in b]
            bx0-=8; by0-=8; bx1+=8; by1+=8
            bx0 = max(bx0, x0); by0 = max(by0, y0); bx1 = min(bx1, x1); by1 = min(by1, y1)
            bg = rgb_bgr(t.get("bg", "#FFFFFF"))
            if (sum(bg) > 3*232 and t.get("mask") != "inpaint") or t.get("mask") == "fill":
                cv2.rectangle(crop, (bx0-x0,by0-y0), (bx1-x0,by1-y0), bg, -1)
            else:
                mask = np.zeros(crop.shape[:2], np.uint8)
                cv2.rectangle(mask, (bx0-x0,by0-y0), (bx1-x0,by1-y0), 255, -1)
                crop = cv2.inpaint(crop, mask, 7, cv2.INPAINT_TELEA)
    for reg in e.get("inpaint", []):
        rx0,ry0,rx1,ry1 = [int(round(v*s)) for v in reg]
        mask = np.zeros(crop.shape[:2], np.uint8); mask[ry0-y0:ry1-y0, rx0-x0:rx1-x0] = 255
        crop = cv2.inpaint(crop, mask, 5, cv2.INPAINT_TELEA)
    path = os.path.join(asset_dir, f"p{i}_{e['id']}.png")
    if e.get("photo"):
        e["transparent"] = False
        path = path[:-4] + ".jpg"; cv2.imwrite(path, crop, [cv2.IMWRITE_JPEG_QUALITY, 90]); return path
    if "transparent" not in e:
        border = np.concatenate([crop[0], crop[-1], crop[:,0], crop[:,-1]])
        e["transparent"] = bool((border.mean(1) > 240).mean() > 0.6 and max(crop.shape[:2]) < 700)
    if e.get("transparent"):
        b, g, r = cv2.split(crop)
        alpha = np.where((b > 235) & (g > 235) & (r > 235), 0, 255).astype(np.uint8)
        crop = cv2.merge([b, g, r, alpha])
    cv2.imwrite(path, crop)
    return path

def rgb_bgr(hexs): hexs = hexs.lstrip("#"); return (int(hexs[4:],16), int(hexs[2:4],16), int(hexs[:2],16))

def set_name(shape, name): shape.name = name

def add_text(slide, e, sc, font, dpi):
    x0,y0,x1,y1 = e["bbox"]
    lines = e["text"].split("\n")
    n = len(lines)
    fs_px = e.get("font_px") or ((y1 - y0) / n) * (0.85 if n == 1 else 0.62)
    if e.get("font_pt"):
        fs_pt = e["font_pt"]; fs_px = fs_pt * dpi / 72 * GLYPH_RATIO
    else:
        # height-based estimate, capped so the longest line fits the OCR width
        by_h = fs_px / GLYPH_RATIO
        longest = max(lines, key=lambda l: measure_width(l, e.get("font", font), e.get("bold")))
        w100 = measure_width(longest, e.get("font", font), e.get("bold"))
        by_w = (x1 - x0) / w100 * 100
        fit = min(by_h, by_w * (1.05 if n == 1 else 1.0)) if len(longest.strip()) >= 3 else by_h
        if fit < 12: print(f"  warn: tiny font for {e.get('id')} {e.get('text','')[:15]!r} fit={fit:.1f}px bbox={e['bbox']}"); fit = max(fit, 12)
        fs_pt = round(fit * 72 / dpi, 1); fs_px = fit * GLYPH_RATIO
    pitch = e.get("pitch_px") or fs_px * 1.2 / GLYPH_RATIO
    # line pitch in px for PPT at spacing m: m * 1.2 * fs_px/GLYPH_RATIO
    mult = e.get("line_spacing") or round(pitch * GLYPH_RATIO / (1.2 * fs_px), 2)
    line_h_px = mult * 1.2 * fs_px / GLYPH_RATIO
    top_off = (line_h_px - fs_px) / 2.0
    pad_x = 0.06 * fs_px + 3
    bx = x0 - pad_x; by = y0 - top_off
    bw = (x1 - x0) + 2*pad_x + (0.6 if n == 1 else 1.2)*fs_px   # slack against wrapping
    bh = n * line_h_px + 2
    tb = slide.shapes.add_textbox(Emu(int(bx*sc)), Emu(int(by*sc)), Emu(int(bw*sc)), Emu(int(bh*sc)))
    set_name(tb, e.get("name", e["id"]))
    tf = tb.text_frame; tf.word_wrap = e.get("wrap", n > 1)
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if e.get("rotation"): tb.rotation = e["rotation"]
    align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}[e.get("align","left")]
    for k, ln in enumerate(lines):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = mult
        runs = e["runs"] if (k == 0 and e.get("runs")) else [dict(text=ln)]
        for ru in runs:
            r = p.add_run(); r.text = ru["text"]
            r.font.size = Pt(fs_pt); r.font.name = e.get("font", font); r.font.bold = bool(ru.get("bold", e.get("bold")))
            r.font.color.rgb = rgb(ru.get("color", e.get("color", "#000000")))
            rPr = r._r.get_or_add_rPr()
            ea = rPr.find(qn("a:ea"))
            if ea is None: ea = etree.SubElement(rPr, qn("a:ea"))
            ea.set("typeface", e.get("font", font))
    if e.get("char_spacing"):
        for p in tf.paragraphs:
            for r in p.runs: r._r.get_or_add_rPr().set("spc", str(int(e["char_spacing"]*100)))
    return tb

def add_rect(slide, e, sc, dpi):
    x0,y0,x1,y1 = e["bbox"]
    geom = MSO_SHAPE.ROUNDED_RECTANGLE if e.get("rounded") else MSO_SHAPE.RECTANGLE
    if e.get("geometry") == "ellipse": geom = MSO_SHAPE.OVAL
    sh = slide.shapes.add_shape(geom, Emu(int(x0*sc)), Emu(int(y0*sc)), Emu(int((x1-x0)*sc)), Emu(int((y1-y0)*sc)))
    set_name(sh, e.get("name", e["id"]))
    if e.get("rounded") and geom == MSO_SHAPE.ROUNDED_RECTANGLE:
        sh.adjustments[0] = e.get("radius", 0.12)
    if e.get("gradient"):
        sh.fill.gradient(); sh.fill.gradient_angle = e.get("gradient_angle", 90)
        st = sh.fill.gradient_stops; st[0].color.rgb = rgb(e["gradient"][0]); st[1].color.rgb = rgb(e["gradient"][1])
    elif e.get("fill"): sh.fill.solid(); sh.fill.fore_color.rgb = rgb(e["fill"])
    else: sh.fill.background()
    if e.get("line"):
        sh.line.color.rgb = rgb(e["line"]); sh.line.width = Pt(max(0.5, e.get("line_px", 1) * 72 / dpi))
    else: sh.line.fill.background()
    if e.get("dash"):
        from pptx.enum.dml import MSO_LINE
        sh.line.dash_style = MSO_LINE.DASH
    sh.shadow.inherit = False
    return sh

def add_line(slide, e, sc, dpi):
    x0,y0,x1,y1 = e["bbox"]
    horiz = (x1-x0) >= (y1-y0)
    if e.get("points"):
        (ax,ay),(bx,by) = e["points"]
    elif horiz: ay = by = (y0+y1)/2; ax, bx = x0, x1
    else: ax = bx = (x0+x1)/2; ay, by = y0, y1
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(ax*sc)), Emu(int(ay*sc)), Emu(int(bx*sc)), Emu(int(by*sc)))
    set_name(c, e.get("name", e["id"]))
    c.line.color.rgb = rgb(e.get("color", "#000000"))
    wpx = e.get("width_px", 1)
    c.line.width = Pt(max(0.5, wpx * 72 / dpi))
    if e.get("dash"):
        from pptx.enum.dml import MSO_LINE
        c.line.dash_style = {"dot": MSO_LINE.ROUND_DOT, "dash": MSO_LINE.DASH}.get(e["dash"], MSO_LINE.DASH)
    ln = c.line._get_or_add_ln()
    if e.get("arrow_end"): etree.SubElement(ln, qn("a:tailEnd")).set("type", "triangle")
    if e.get("arrow_start"): etree.SubElement(ln, qn("a:headEnd")).set("type", "triangle")
    return c

TABLE_NO_STYLE = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"   # No Style, No Grid

def _image_boxes(layout):
    """실제로 슬라이드에 그려질 이미지들의 bbox (표 셀 배경 판단용)"""
    skip = layout.get("_skip_imgs", set())
    return [el["bbox"] for el in layout["elements"]
            if el["type"] == "image" and el.get("id") not in skip]

def add_table(slide, work, layout, e, sc, font, dpi):
    """표 요소 -> 실제 PowerPoint 표 개체"""
    import tables as tbl
    if "cells" not in e:
        img = cv2.imread(f"{work}/pages/p{layout['page']}.png")
        t = tbl.build_table(img, tbl.ocr_lines(work, layout["page"]), e["bbox"], e.get("xs"), e.get("ys"))
        for k, v in t.items(): e.setdefault(k, v)
    xs, ys, cells = e["xs"], e["ys"], e["cells"]
    rows, cols = len(ys) - 1, len(xs) - 1
    for pos, txt in (e.get("fix") or {}).items():          # overrides: {"1,2": "교정문구"}
        r, c = [int(v) for v in pos.split(",")]
        if 0 <= r < rows and 0 <= c < cols: cells[r][c] = txt
    gf = slide.shapes.add_table(rows, cols, Emu(int(xs[0]*sc)), Emu(int(ys[0]*sc)),
                                Emu(int((xs[-1]-xs[0])*sc)), Emu(int((ys[-1]-ys[0])*sc)))
    shape = gf; table = gf.table
    set_name(shape, e.get("name", e["id"]))
    tblPr = table._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        tblPr.set("firstRow", "0"); tblPr.set("bandRow", "0")
        st = tblPr.find(qn("a:tableStyleId"))
        if st is None: st = etree.SubElement(tblPr, qn("a:tableStyleId"))
        st.text = e.get("style_id", TABLE_NO_STYLE)
    for c in range(cols): table.columns[c].width = Emu(int((xs[c+1]-xs[c])*sc))
    for r in range(rows): table.rows[r].height = Emu(int((ys[r+1]-ys[r])*sc))
    line_col = e.get("line", "#B4B8B9"); line_pt = e.get("line_px", 2) * 72 / dpi
    imgs = _image_boxes(layout)
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.margin_left = cell.margin_right = Emu(int(4*sc))
            cell.margin_top = cell.margin_bottom = Emu(int(2*sc))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            fill = (e.get("fills") or [[None]*cols]*rows)[r][c]
            cb = [xs[c], ys[r], xs[c+1], ys[r+1]]
            area = max(1, (cb[2]-cb[0]) * (cb[3]-cb[1]))
            for ib in imgs:                       # 셀의 절반 이상을 이미지가 덮으면 채우기 없음
                ov = max(0, min(cb[2], ib[2]) - max(cb[0], ib[0])) * max(0, min(cb[3], ib[3]) - max(cb[1], ib[1]))
                if ov / area > 0.5: fill = None; break
            if fill and fill.upper() not in ("#FFFFFF", "#FEFEFE", "#FDFDFD", "#FCFCFC"):
                cell.fill.solid(); cell.fill.fore_color.rgb = rgb(fill)
            else:
                cell.fill.background()
            st = ((e.get("styles") or [[None]*cols]*rows)[r][c]) or {}
            txt = cells[r][c]
            tf = cell.text_frame; tf.word_wrap = True
            fs_pt = st.get("font_pt")
            if not fs_pt:
                fit = st.get("font_px", 26) / GLYPH_RATIO
                avail = (xs[c+1] - xs[c]) - 18
                lines_ = [ln for ln in txt.split(chr(10)) if ln.strip()]
                if lines_ and avail > 20:
                    longest = max(lines_, key=lambda ln: measure_width(ln, e.get("font", font), st.get("bold")))
                    w100 = measure_width(longest, e.get("font", font), st.get("bold"))
                    fit = min(fit, avail / w100 * 98)
                fs_pt = round(max(6.0, fit * 72 / dpi), 1)
            for k, ln in enumerate(txt.split("\n")):
                p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                               "right": PP_ALIGN.RIGHT}[st.get("align", "left")]
                run = p.add_run(); run.text = ln
                run.font.size = Pt(fs_pt); run.font.name = e.get("font", font)
                run.font.bold = bool(st.get("bold")); run.font.color.rgb = rgb(st.get("color", "#3A3C3C"))
                rPr = run._r.get_or_add_rPr()
                ea = rPr.find(qn("a:ea"))
                if ea is None: ea = etree.SubElement(rPr, qn("a:ea"))
                ea.set("typeface", e.get("font", font))
            _cell_borders(cell, line_col, line_pt, e.get("borders", "all"))
    return shape

def _cell_borders(cell, color, width_pt, mode="all"):
    tcPr = cell._tc.get_or_add_tcPr()
    tags = {"L": "a:lnL", "R": "a:lnR", "T": "a:lnT", "B": "a:lnB"}
    want = "LRTB" if mode == "all" else mode
    for k, tag in tags.items():
        el = tcPr.find(qn(tag))
        if el is not None: tcPr.remove(el)
        if k not in want: continue
        ln = etree.SubElement(tcPr, qn(tag))
        ln.set("w", str(int(width_pt * 12700))); ln.set("cap", "flat"); ln.set("cmpd", "sng"); ln.set("algn", "ctr")
        fill = etree.SubElement(ln, qn("a:solidFill"))
        clr = etree.SubElement(fill, qn("a:srgbClr")); clr.set("val", color.lstrip("#").upper())
    # 순서 보정: lnL, lnR, lnT, lnB 는 tcPr 자식 중 앞쪽에 와야 한다
    order = [qn(t) for t in ("a:lnL", "a:lnR", "a:lnT", "a:lnB", "a:lnTlToBr", "a:lnBlToTr")]
    kids = list(tcPr)
    kids.sort(key=lambda el: order.index(el.tag) if el.tag in order else len(order))
    for el in kids: tcPr.append(el)

def add_image(slide, work, layout, e, sc, asset_dir):
    path = e.get("asset")
    if path and not os.path.isabs(path): path = os.path.join(asset_dir, path)
    if not path or not os.path.exists(path): path = crop_image(work, layout, e, asset_dir)
    x0,y0,x1,y1 = e["bbox"]
    pic = slide.shapes.add_picture(path, Emu(int(x0*sc)), Emu(int(y0*sc)), Emu(int((x1-x0)*sc)), Emu(int((y1-y0)*sc)))
    set_name(pic, e.get("name", e["id"]))
    return pic

def main():
    work, out = sys.argv[1], sys.argv[2]
    pages = None; font = os.environ.get("PDF2PPTX_FONT", "맑은 고딕")
    if "--pages" in sys.argv: pages = [int(v) for v in sys.argv[sys.argv.index("--pages")+1].split(",")]
    if "--font" in sys.argv: font = sys.argv[sys.argv.index("--font")+1]
    files = sorted(glob.glob(f"{work}/layout/p*.json"), key=lambda f: int(os.path.basename(f)[1:-5]))
    asset_dir = f"{work}/assets"; os.makedirs(asset_dir, exist_ok=True)
    prs = None; counts = []
    for f in files:
        layout = json.load(open(f, encoding="utf-8"))
        i = layout["page"]
        if pages and i not in pages: continue
        ovf = f"{work}/overrides/p{i}.json"
        if os.path.exists(ovf): apply_overrides(layout, json.load(open(ovf, encoding="utf-8")))
        if prs is None:
            prs = Presentation()
            prs.slide_width = Emu(int(layout["pt_w"] * 12700)); prs.slide_height = Emu(int(layout["pt_h"] * 12700))
        sc = prs.slide_width / layout["width"]   # EMU per px
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        order = {"rect": 0, "table": 1, "image": 2, "line": 3, "text": 4}
        els = sorted(layout["elements"], key=lambda e: (e.get("z", order[e["type"]])))
        c = dict(text=0, rect=0, line=0, image=0, table=0)
        tbls = [t for t in els if t["type"] == "table"]
        def _owner(el):
            b = el["bbox"]; cx, cy = (b[0]+b[2])/2, (b[1]+b[3])/2
            for t in tbls:
                if t["bbox"][0]-4 <= cx <= t["bbox"][2]+4 and t["bbox"][1]-4 <= cy <= t["bbox"][3]+4: return t
            return None
        def in_table(el): return _owner(el) is not None
        def hides_text(el):
            """표 안의 이미지가 글자 있는 셀을 덮는지"""
            t = _owner(el)
            if not t: return False
            b = el["bbox"]; xs, ys = t["xs"], t["ys"]
            for r in range(t["rows"]):
                for c in range(t["cols"]):
                    if not t["cells"][r][c].strip(): continue
                    ov = (max(0, min(b[2], xs[c+1]) - max(b[0], xs[c]))
                          * max(0, min(b[3], ys[r+1]) - max(b[1], ys[r])))
                    if ov > 0.4 * max(1, (xs[c+1]-xs[c]) * (ys[r+1]-ys[r])): return True
            return False
        layout["_skip_imgs"] = {e.get("id") for e in els if e["type"] == "image" and in_table(e)
                                and (not e.get("photo") or hides_text(e))}
        for e in els:
            if e["type"] in ("text", "line") and in_table(e): continue
            # 표 안의 비사진 이미지는 셀 배경/글자가 이미 표로 복원되므로 제외
            if e["type"] == "image" and e.get("id") in layout["_skip_imgs"]: continue
            if e["type"] == "text": add_text(slide, e, sc, font, layout["dpi"])
            elif e["type"] == "rect": add_rect(slide, e, sc, layout["dpi"])
            elif e["type"] == "line": add_line(slide, e, sc, layout["dpi"])
            elif e["type"] == "image": add_image(slide, work, layout, e, sc, asset_dir)
            elif e["type"] == "table": add_table(slide, work, layout, e, sc, font, layout["dpi"])
            c[e["type"]] += 1
        counts.append((i, c))
    prs.save(out)
    for i, c in counts: print(f"slide {i}: {c}")
    print("saved", out)

if __name__ == "__main__":
    main()
