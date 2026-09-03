# -*- coding: utf-8 -*-
"""레이아웃 -> PPTX 빌드와 패키지 검증.

여기서 확인하는 것: 실제 편집 가능한 개체가 나오는가, 겹침 순서가 의도대로인가,
그리고 손상된 패키지를 검증기가 실제로 거부하는가(음성 대조군).
"""
import os, sys, json, zipfile, shutil
import pytest

import build as build_mod
import quality


def _build(work, out, pages=None):
    argv = sys.argv
    sys.argv = ["build", work, out] + (["--pages", pages] if pages else [])
    try:
        build_mod.main()
    finally:
        sys.argv = argv
    return out


@pytest.fixture
def pptx(work, tmp_path):
    return _build(work, str(tmp_path / "out.pptx"))


# --------------------------------------------------------------- 편집 가능 개체
def test_build_creates_native_objects_not_one_big_picture(pptx):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(pptx)
    assert len(prs.slides) == 1
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) >= 4, "개체가 통짜 이미지 하나로 뭉개졌다"

    texts = [s for s in shapes if s.has_text_frame and s.text_frame.text.strip()]
    pics = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert [s.text_frame.text for s in texts].count("보고서 제목") == 1
    assert len(pics) == 1, "사진은 개별 그림 개체여야 한다"
    # 슬라이드를 통째로 덮는 그림이 있으면 편집 가능성이 사라진 것이다
    for p in pics:
        assert not (p.width >= prs.slide_width * 0.98 and p.height >= prs.slide_height * 0.98)


def test_text_is_editable_and_keeps_styling(pptx):
    from pptx import Presentation
    prs = Presentation(pptx)
    title = next(s for s in prs.slides[0].shapes
                 if s.has_text_frame and s.text_frame.text.strip() == "보고서 제목")
    run = title.text_frame.paragraphs[0].runs[0]
    assert run.font.bold is True
    assert run.font.size.pt > 0
    body = next(s for s in prs.slides[0].shapes
                if s.has_text_frame and "첫째 줄" in s.text_frame.text)
    assert len(body.text_frame.paragraphs) == 2, "여러 줄 문단이 한 줄로 뭉치면 안 된다"


def test_slide_size_matches_source_page(work, pptx):
    from pptx import Presentation
    lay = json.load(open(os.path.join(work, "layout", "p1.json"), encoding="utf-8"))
    prs = Presentation(pptx)
    assert prs.slide_width == int(lay["pt_w"] * 12700)
    assert prs.slide_height == int(lay["pt_h"] * 12700)


# --------------------------------------------------------------- 겹침 순서
def test_background_rect_is_placed_under_its_text(pptx):
    """현재는 타입 기반 순서(rect < text). 원본 paint order 복원은 ROADMAP R2.

    이 테스트는 '적어도 배경이 글자를 덮지는 않는다'는 최소 보장을 고정한다.
    """
    from pptx import Presentation
    prs = Presentation(pptx)
    shapes = list(prs.slides[0].shapes)
    idx = {}
    for i, s in enumerate(shapes):
        if s.has_text_frame and s.text_frame.text.strip() == "보고서 제목":
            idx["text"] = i
        if s.shape_type is not None and not s.has_text_frame:
            idx.setdefault("first_non_text", i)
    assert "text" in idx
    assert idx["text"] > idx.get("first_non_text", -1)


def test_explicit_z_override_wins_over_type_order(work, tmp_path):
    """overrides 로 z 를 주면 타입 기본 순서를 이길 수 있어야 한다(paint order 복원의 토대)."""
    from pptx import Presentation
    p = os.path.join(work, "layout", "p1.json")
    lay = json.load(open(p, encoding="utf-8"))
    for e in lay["elements"]:
        if e["id"] == "t0": e["z"] = -5       # 제목을 맨 아래로
    json.dump(lay, open(p, "w", encoding="utf-8"), ensure_ascii=False)
    prs = Presentation(_build(work, str(tmp_path / "z.pptx")))
    first = list(prs.slides[0].shapes)[0]
    assert first.has_text_frame and first.text_frame.text.strip() == "보고서 제목"


# --------------------------------------------------------------- overrides
def test_overrides_delete_and_set(work, tmp_path):
    from pptx import Presentation
    os.makedirs(os.path.join(work, "overrides"), exist_ok=True)
    json.dump({"delete": ["l0"], "set": {"t0": {"text": "교정된 제목"}}},
              open(os.path.join(work, "overrides", "p1.json"), "w", encoding="utf-8"),
              ensure_ascii=False)
    prs = Presentation(_build(work, str(tmp_path / "ov.pptx")))
    texts = [s.text_frame.text for s in prs.slides[0].shapes if s.has_text_frame]
    assert "교정된 제목" in texts and "보고서 제목" not in texts


# --------------------------------------------------------------- 패키지 검증
def test_valid_package_passes(pptx):
    assert quality.validate_package(pptx) is None


def test_truncated_file_is_rejected(pptx, tmp_path):
    bad = tmp_path / "truncated.pptx"
    data = open(pptx, "rb").read()
    open(bad, "wb").write(data[:len(data) // 2])
    assert quality.validate_package(str(bad)) is not None


def test_malformed_xml_is_rejected(pptx, tmp_path):
    """zip 은 멀쩡하지만 XML 이 깨진 경우 -- 실제로 PowerPoint 복구 경고가 뜨는 상황."""
    bad = tmp_path / "badxml.pptx"
    with zipfile.ZipFile(pptx) as src, zipfile.ZipFile(bad, "w") as dst:
        for it in src.infolist():
            data = src.read(it.filename)
            if it.filename == "ppt/presentation.xml":
                data = data[:-10] + b"<unclosed>"
            dst.writestr(it, data)
    assert quality.validate_package(str(bad)) is not None


def test_missing_required_part_is_rejected(pptx, tmp_path):
    bad = tmp_path / "nopart.pptx"
    with zipfile.ZipFile(pptx) as src, zipfile.ZipFile(bad, "w") as dst:
        for it in src.infolist():
            if it.filename == "[Content_Types].xml":
                continue
            dst.writestr(it, src.read(it.filename))
    assert "[Content_Types].xml" in (quality.validate_package(str(bad)) or "")
