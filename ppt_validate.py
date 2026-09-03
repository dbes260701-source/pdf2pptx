# -*- coding: utf-8 -*-
"""
ppt_validate.py -- Microsoft PowerPoint 호환성 검증 게이트 (ROADMAP R4)

평가서 [B] §4.8 대응.

  "패키지 검증은 PowerPoint 호환성과 같지 않다. zip/XML 이 온전하고 python-pptx 로
   다시 열린다는 것이, PowerPoint 가 복구 경고 없이 연다는 뜻은 아니다."

접근 방식에 대해:
PowerPoint 가 파일을 '복구'했는지는 COM 으로 직접 알려주지 않는다. 복구는 대화상자로
알리는데, 자동화에서는 그 대화상자를 신뢰성 있게 잡을 수 없다.
그래서 대화상자를 엿보는 대신 **왕복 후 내용을 대조**한다.

  1. python-pptx 로 원본 개체 목록을 뜬다
  2. PowerPoint 로 열고 -> 저장하고 -> 닫는다 (PowerPoint 자신의 직렬화로 다시 씀)
  3. 저장된 파일의 개체 목록을 다시 뜬다
  4. 둘을 비교한다

PowerPoint 가 무언가를 버렸거나 고쳐 썼다면 개체 수·종류·글자가 달라진다.
"열렸다"가 아니라 "열고 저장해도 내용이 그대로다"를 확인하는 것이므로,
복구 여부보다 오히려 강한 검사다.

PowerPoint 가 없는 환경(LibreOffice 전용, CI)에서는 status="unavailable" 로 보고하고
실패로 만들지 않는다. 다만 그 경우 **PowerPoint 호환을 주장하지 않는다**.
"""
import os


def available():
    """PowerPoint COM 을 쓸 수 있는지."""
    try:
        import win32com.client, pythoncom      # noqa: F401
    except ImportError:
        return False
    try:
        pythoncom.CoInitialize()
        app = win32com.client.Dispatch("PowerPoint.Application")
        try:
            app.Quit()
        except Exception:
            pass
        return True
    except Exception:
        return False


def inventory(pptx):
    """비교용 개체 목록. PowerPoint 가 내용을 바꿨는지 보기 위한 최소 지문."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(pptx)
    slides = []
    for s in prs.slides:
        shapes, texts, pics, tables = 0, [], 0, []
        for sh in s.shapes:
            shapes += 1
            if sh.has_text_frame and sh.text_frame.text.strip():
                texts.append(sh.text_frame.text.strip())
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pics += 1
            if getattr(sh, "has_table", False) and sh.has_table:
                tables.append((len(sh.table.rows), len(sh.table.columns)))
        slides.append(dict(shapes=shapes, texts=sorted(texts), pictures=pics,
                           tables=sorted(tables)))
    return dict(slide_count=len(prs.slides),
                width=int(prs.slide_width), height=int(prs.slide_height),
                slides=slides)


# PowerPoint 는 슬라이드 크기를 자신의 내부 정밀도로 반올림해 다시 쓴다.
# 실측: 7563611x10693907 -> 7562850x10693400 EMU (0.06pt 차이). 손실이 아니라 반올림이므로
# 1pt(12700 EMU) 까지는 같은 크기로 본다. 용지 자체가 바뀌면 이보다 훨씬 크게 벌어진다.
SIZE_TOLERANCE_EMU = 12700


def _diff(before, after):
    """왕복 전후 차이를 사람이 읽을 수 있는 사유 목록으로."""
    out = []
    if before["slide_count"] != after["slide_count"]:
        out.append(f"슬라이드 수 {before['slide_count']} -> {after['slide_count']}")
        return out                       # 장수가 다르면 이하 비교는 의미 없다
    dw = abs(before["width"] - after["width"])
    dh = abs(before["height"] - after["height"])
    if dw > SIZE_TOLERANCE_EMU or dh > SIZE_TOLERANCE_EMU:
        out.append(f"슬라이드 크기 {before['width']}x{before['height']} -> "
                   f"{after['width']}x{after['height']}")
    for i, (b, a) in enumerate(zip(before["slides"], after["slides"]), start=1):
        if b["shapes"] != a["shapes"]:
            out.append(f"{i}번 슬라이드 개체 수 {b['shapes']} -> {a['shapes']}")
        if b["pictures"] != a["pictures"]:
            out.append(f"{i}번 슬라이드 그림 {b['pictures']} -> {a['pictures']}")
        if b["tables"] != a["tables"]:
            out.append(f"{i}번 슬라이드 표 {b['tables']} -> {a['tables']}")
        lost = [t for t in b["texts"] if t not in a["texts"]]
        if lost:
            sample = "; ".join(t[:20] for t in lost[:3])
            out.append(f"{i}번 슬라이드 글자 {len(lost)}건 유실 ({sample})")
    return out


def roundtrip(pptx, outdir):
    """PowerPoint 로 열고 -> 저장하고 -> 다시 열어 내용이 보존되는지.

    반환: dict(status, differences, saved, opened, error, ...)
      status = "pass" | "fail" | "unavailable"
    """
    result = dict(status="unavailable", opened=False, saved=False,
                  differences=[], error=None, roundtrip_path=None)
    try:
        import win32com.client, pythoncom
    except ImportError:
        result["error"] = "pywin32 없음"
        return result

    try:
        before = inventory(pptx)
    except Exception as ex:
        result["status"] = "fail"
        result["error"] = f"원본을 python-pptx 로 열 수 없음: {ex}"
        return result

    os.makedirs(outdir, exist_ok=True)
    saved = os.path.abspath(os.path.join(outdir, "roundtrip.pptx"))
    if os.path.exists(saved):
        os.remove(saved)

    app = None
    pythoncom.CoInitialize()
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        try:
            app.DisplayAlerts = 1          # ppAlertsNone -- 대화상자로 멈추지 않게
        except Exception:
            pass
        # 읽기 전용으로 열지 않는다. 저장까지 해봐야 PowerPoint 의 직렬화를 검증한다.
        pres = app.Presentations.Open(os.path.abspath(pptx), ReadOnly=False,
                                      Untitled=False, WithWindow=False)
        result["opened"] = True
        pres.SaveCopyAs(saved)             # 원본을 건드리지 않고 PowerPoint 형식으로 다시 씀
        pres.Close()
        result["saved"] = os.path.exists(saved)
    except Exception as ex:
        result["status"] = "fail"
        result["error"] = f"PowerPoint 열기/저장 실패: {ex}"
        return result
    finally:
        if app is not None:
            try:
                app.Quit()
            except Exception:
                pass

    if not result["saved"]:
        result["status"] = "fail"
        result["error"] = "PowerPoint 가 저장본을 만들지 못했습니다."
        return result

    result["roundtrip_path"] = saved
    try:
        after = inventory(saved)
    except Exception as ex:
        result["status"] = "fail"
        result["error"] = f"PowerPoint 저장본을 다시 열 수 없음: {ex}"
        return result

    result["differences"] = _diff(before, after)
    result["status"] = "fail" if result["differences"] else "pass"
    return result
